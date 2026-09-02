from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

import fitz
from docx import Document as WordDocument
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx"}


class DocumentExtractionError(Exception):
    """Raised when an uploaded document cannot be parsed."""


@dataclass(frozen=True)
class ExtractedSection:
    page_number: int
    section_label: str
    text_content: str


@dataclass(frozen=True)
class ExtractionResult:
    sections: list[ExtractedSection]
    page_count: int


def _clean_text(parts: list[str]) -> str:
    cleaned = [line.strip() for part in parts for line in part.splitlines() if line.strip()]
    return "\n".join(cleaned)


def _extract_pdf(content: bytes) -> ExtractionResult:
    sections: list[ExtractedSection] = []
    with fitz.open(stream=content, filetype="pdf") as document:
        for index, page in enumerate(document, start=1):
            sections.append(
                ExtractedSection(index, f"Halaman {index}", _clean_text([page.get_text("text")]))
            )
        return ExtractionResult(sections=sections, page_count=document.page_count)


def _shape_text(shape: object) -> list[str]:
    parts: list[str] = []
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            parts.extend(_shape_text(child))
        return parts
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells if cell.text.strip()))
    elif getattr(shape, "has_text_frame", False):
        parts.append(shape.text)
    return parts


def _extract_pptx(content: bytes) -> ExtractionResult:
    presentation = Presentation(BytesIO(content))
    sections: list[ExtractedSection] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            parts.extend(_shape_text(shape))
        sections.append(ExtractedSection(index, f"Slide {index}", _clean_text(parts)))
    return ExtractionResult(sections=sections, page_count=len(presentation.slides))


def _extract_docx(content: bytes) -> ExtractionResult:
    document = WordDocument(BytesIO(content))
    parts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells if cell.text.strip()))
    text = _clean_text(parts)
    return ExtractionResult(
        sections=[ExtractedSection(1, "Dokumen", text)],
        page_count=1,
    )


def extract_document(content: bytes, file_name: str) -> ExtractionResult:
    extension = Path(file_name).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise DocumentExtractionError("Format file harus PDF, PPTX, atau DOCX.")
    try:
        if extension == ".pdf":
            return _extract_pdf(content)
        if extension == ".pptx":
            return _extract_pptx(content)
        return _extract_docx(content)
    except DocumentExtractionError:
        raise
    except Exception as exc:
        raise DocumentExtractionError(
            "File tidak dapat dibaca. Pastikan dokumen tidak rusak atau diproteksi password."
        ) from exc
