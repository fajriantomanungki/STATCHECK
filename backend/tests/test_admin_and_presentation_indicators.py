from io import BytesIO

from pptx import Presentation


def auth_headers(client):
    response = client.post("/api/v1/auth/login", json={"nik": "admin", "password": "Admin123!"})
    return {"Authorization": f"Bearer {response.json()['access_token']}"}


def pptx_bytes(text: str) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Indikator Pariwisata"
    slide.placeholders[1].text = text
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def create_brs(client, headers, name: str = "Perkembangan Pariwisata Juli 2026") -> str:
    response = client.post("/api/v1/brs", headers=headers, json={
        "nama_brs": name, "waktu_rilis": "2026-09-01",
        "fungsi_pj": "Statistik Distribusi", "supervisor_id": None,
        "team_user_ids": [],
    })
    assert response.status_code == 201, response.text
    return response.json()["id"]


def upload_presentation(client, headers, brs_id: str, text: str):
    response = client.post(
        f"/api/v1/brs/{brs_id}/documents", headers=headers,
        data={"document_type": "bahan_paparan"},
        files={"file": ("paparan.pptx", pptx_bytes(text))},
    )
    assert response.status_code == 201, response.text
    return client.get(
        f"/api/v1/brs/{brs_id}/presentation-indicators", headers=headers
    ).json()


def test_admin_can_create_edit_and_delete_user(client):
    headers = auth_headers(client)
    created = client.post("/api/v1/users", headers=headers, json={
        "nama": "PJK Baru", "nik": "pjk-baru", "user_level": "pjk",
        "fungsi": "Statistik Distribusi", "password": "PjkBaru123!",
    })
    assert created.status_code == 201, created.text
    user_id = created.json()["id"]

    updated = client.put(f"/api/v1/users/{user_id}", headers=headers, json={
        "nama": "PJK Diperbarui", "nik": "pjk-baru", "user_level": "pjk",
        "fungsi": "Statistik Sosial", "is_active": False, "password": None,
    })
    assert updated.status_code == 200, updated.text
    assert updated.json()["is_active"] is False

    deleted = client.delete(f"/api/v1/users/{user_id}", headers=headers)
    assert deleted.status_code == 204, deleted.text


def test_admin_can_delete_unused_indicator(client):
    headers = auth_headers(client)
    created = client.post("/api/v1/indicators", headers=headers, json={
        "nama_indikator": "Indikator Sementara", "kategori": "Uji",
        "satuan_default": "unit", "fungsi": "Pengujian",
    })
    assert created.status_code == 201
    deleted = client.delete(f"/api/v1/indicators/{created.json()['id']}", headers=headers)
    assert deleted.status_code == 204


def test_presentation_upload_extracts_editable_indicator_table(client):
    headers = auth_headers(client)
    brs_id = create_brs(client, headers)
    text = (
        "TPK hotel bintang pada Juli 2026 tercatat 52,31 persen. "
        "Rata-rata Lama Menginap Tamu hotel bintang pada Juli 2026 tercatat 1,67 hari."
    )
    upload = client.post(
        f"/api/v1/brs/{brs_id}/documents", headers=headers,
        data={"document_type": "bahan_paparan"},
        files={"file": ("paparan.pptx", pptx_bytes(text))},
    )
    assert upload.status_code == 201, upload.text

    listed = client.get(f"/api/v1/brs/{brs_id}/presentation-indicators", headers=headers)
    assert listed.status_code == 200
    rows = listed.json()
    assert len(rows) >= 2
    tpk = next(row for row in rows if "TPK" in row["indicator_name"])
    assert tpk["value_text"] == "52,31"
    assert tpk["period_label"] == "Juli 2026"
    assert tpk["data_type"] == "percentage"
    assert "52,31" in tpk["metadata_text"]

    saved = client.put(
        f"/api/v1/brs/{brs_id}/presentation-indicators/{tpk['id']}",
        headers=headers,
        json={"analysis": "TPK menunjukkan peningkatan.", "phenomenon": "Aktivitas perjalanan meningkat."},
    )
    assert saved.status_code == 200
    assert saved.json()["analysis"] == "TPK menunjukkan peningkatan."

    refreshed = client.post(
        f"/api/v1/brs/{brs_id}/presentation-indicators/refresh", headers=headers
    )
    assert refreshed.status_code == 200
    refreshed_tpk = next(row for row in refreshed.json() if "TPK" in row["indicator_name"])
    assert refreshed_tpk["analysis"] == "TPK menunjukkan peningkatan."


def test_presentation_extraction_ignores_ambiguous_and_non_actual_numbers(client):
    headers = auth_headers(client)
    brs_id = create_brs(client, headers, "BRS Penyaringan Angka")
    rows = upload_presentation(client, headers, brs_id, (
        "Bahan paparan ini terdiri dari 12 slide. "
        "Target TPK tahun 2027 sebesar 60 persen. "
        "TPK hotel bintang pada Juli 2026 tercatat 52,31 persen. "
        "Nomor dokumen 49/08/72."
    ))

    assert [row["value_text"] for row in rows] == ["52,31"]
    assert rows[0]["indicator_name"] == "Tingkat Penghunian Kamar (TPK) — Hotel Bintang"


def test_presentation_extraction_understands_total_and_guest_components(client):
    headers = auth_headers(client)
    brs_id = create_brs(client, headers, "BRS Komponen Tamu")
    rows = upload_presentation(client, headers, brs_id, (
        "Jumlah tamu yang menginap pada akomodasi Hotel Bintang pada Juli 2026 "
        "tercatat 27.073 orang, terdiri dari 26.621 orang tamu domestik dan "
        "452 orang tamu mancanegara."
    ))

    assert {row["value_text"] for row in rows} == {"27.073", "26.621", "452"}
    names = {row["value_text"]: row["indicator_name"] for row in rows}
    assert names["27.073"].endswith("Hotel Bintang — Total")
    assert names["26.621"].endswith("Hotel Bintang — Domestik")
    assert names["452"].endswith("Hotel Bintang — Mancanegara")


def test_presentation_extraction_distinguishes_level_change_and_range(client):
    headers = auth_headers(client)
    brs_id = create_brs(client, headers, "BRS Peran Nilai")
    rows = upload_presentation(client, headers, brs_id, (
        "TPK Hotel Bintang pada Juli 2026 tercatat 51,95 persen, naik 0,48 "
        "persen poin dibandingkan Juni 2026. Rata-rata setiap tamu Hotel Bintang "
        "pada Juli 2026 menghabiskan waktu sekitar 1 hingga 2 hari untuk menginap."
    ))

    by_value = {row["value_text"]: row for row in rows}
    assert by_value["51,95"]["value_role"] == "level"
    assert by_value["0,48"]["value_role"] == "change"
    assert by_value["0,48"]["indicator_name"].endswith("Perubahan")
    assert by_value["1 hingga 2"]["data_type"] == "range"
    assert by_value["1 hingga 2"]["indicator_name"].startswith("Rata-rata Lama Menginap Tamu")
