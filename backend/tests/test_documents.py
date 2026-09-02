from io import BytesIO

import fitz
from docx import Document as WordDocument
from fastapi.testclient import TestClient
from pptx import Presentation

from app.services.document_extractor import extract_document


def auth_headers(client: TestClient) -> dict[str, str]:
    response = client.post("/api/v1/auth/login", json={"nik": "admin", "password": "Admin123!"})
    assert response.status_code == 200
    return {"Authorization": f"Bearer {response.json()['access_token']}"}


def create_brs(client: TestClient, headers: dict[str, str]) -> str:
    response = client.post(
        "/api/v1/brs",
        headers=headers,
        json={
            "nama_brs": "Perkembangan Pariwisata September 2026",
            "waktu_rilis": "2026-10-01",
            "fungsi_pj": "Statistik Distribusi",
            "supervisor_id": None,
            "team_user_ids": [],
        },
    )
    assert response.status_code == 201
    return response.json()["id"]


def pdf_bytes(text: str) -> bytes:
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), text)
    content = document.tobytes()
    document.close()
    return content


def pptx_bytes(text: str) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Bahan Paparan"
    slide.placeholders[1].text = text
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def docx_bytes(text: str) -> bytes:
    document = WordDocument()
    document.add_heading("Narasi Pimpinan", level=1)
    document.add_paragraph(text)
    stream = BytesIO()
    document.save(stream)
    return stream.getvalue()


def upload(
    client: TestClient,
    headers: dict[str, str],
    brs_id: str,
    document_type: str,
    file_name: str,
    content: bytes,
):
    return client.post(
        f"/api/v1/brs/{brs_id}/documents",
        headers=headers,
        data={"document_type": document_type},
        files={"file": (file_name, content)},
    )


def test_extract_supported_document_formats():
    pdf = extract_document(pdf_bytes("Wisnus 1.007,74 ribu perjalanan"), "publikasi.pdf")
    pptx = extract_document(pptx_bytes("TPK mencapai 52,31 persen"), "paparan.pptx")
    docx = extract_document(docx_bytes("Inflasi sebesar 0,90 persen"), "narasi.docx")

    assert pdf.page_count == 1
    assert "1.007,74" in pdf.sections[0].text_content
    assert pptx.page_count == 1
    assert "52,31" in pptx.sections[0].text_content
    assert docx.page_count == 1
    assert "0,90" in docx.sections[0].text_content


def test_upload_three_documents_and_download(client: TestClient):
    headers = auth_headers(client)
    brs_id = create_brs(client, headers)

    files = [
        ("bahan_publikasi", "publikasi.pdf", pdf_bytes("Publikasi Pariwisata 1.007,74")),
        ("bahan_paparan", "paparan.pptx", pptx_bytes("Paparan Pariwisata 1.007,74")),
        ("narasi_pimpinan", "narasi.docx", docx_bytes("Narasi Pariwisata 1.007,74")),
    ]
    uploaded = []
    for document_type, name, content in files:
        response = upload(client, headers, brs_id, document_type, name, content)
        assert response.status_code == 201, response.text
        payload = response.json()
        assert payload["extraction_status"] == "completed"
        assert payload["page_count"] == 1
        assert payload["extracted_char_count"] > 0
        uploaded.append((payload, content))

    brs = client.get(f"/api/v1/brs/{brs_id}", headers=headers)
    assert brs.status_code == 200
    assert brs.json()["status"] == "documents_uploaded"
    assert brs.json()["jumlah_dokumen"] == 3

    listed = client.get(f"/api/v1/brs/{brs_id}/documents", headers=headers)
    assert listed.status_code == 200
    assert len(listed.json()) == 3

    document, original_content = uploaded[0]
    detail = client.get(f"/api/v1/documents/{document['id']}", headers=headers)
    assert detail.status_code == 200
    assert detail.json()["contents"][0]["text_content"]

    reextract = client.post(f"/api/v1/documents/{document['id']}/reextract", headers=headers)
    assert reextract.status_code == 200
    assert reextract.json()["extraction_status"] == "completed"
    assert reextract.json()["contents"][0]["text_content"]

    download = client.get(f"/api/v1/documents/{document['id']}/download", headers=headers)
    assert download.status_code == 200
    assert download.content == original_content


def test_upload_creates_version_history(client: TestClient):
    headers = auth_headers(client)
    brs_id = create_brs(client, headers)

    first = upload(
        client, headers, brs_id, "bahan_publikasi", "publikasi-v1.pdf", pdf_bytes("Versi pertama")
    )
    second = upload(
        client, headers, brs_id, "bahan_publikasi", "publikasi-v2.pdf", pdf_bytes("Versi kedua")
    )
    assert first.status_code == 201
    assert second.status_code == 201
    assert first.json()["version"] == 1
    assert second.json()["version"] == 2

    active = client.get(f"/api/v1/brs/{brs_id}/documents", headers=headers).json()
    history = client.get(
        f"/api/v1/brs/{brs_id}/documents?include_archived=true", headers=headers
    ).json()
    assert len(active) == 1
    assert active[0]["version"] == 2
    assert len(history) == 2
    assert {item["status"] for item in history} == {"active", "archived"}


def test_rejects_unsupported_file(client: TestClient):
    headers = auth_headers(client)
    brs_id = create_brs(client, headers)
    response = upload(client, headers, brs_id, "bahan_publikasi", "catatan.txt", b"teks")
    assert response.status_code == 415
